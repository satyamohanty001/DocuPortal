"""
Enhanced document loaders supporting multiple formats with table and image extraction.
Supports: PDF, DOCX, TXT, MD, PPT, XLSX, CSV, and SQL databases.
"""

import os
import io
import base64
import sqlite3
import pandas as pd
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
from PIL import Image
import pytesseract
from pdf2image import convert_from_path

# LangChain imports
from langchain.schema import Document
from langchain_community.document_loaders import (
    PyPDFLoader, 
    Docx2txtLoader, 
    TextLoader,
    CSVLoader
)

# Document processing imports
import fitz  # PyMuPDF
from pptx import Presentation
from openpyxl import load_workbook
import markdown
from sqlalchemy import create_engine, text
import pymongo

from logger import GLOBAL_LOGGER as log
from exception.custom_exception import DocumentPortalException

class EnhancedDocumentLoader:
    """
    Enhanced document loader supporting multiple formats with advanced features:
    - Table extraction from PDFs, Excel, PowerPoint
    - Image extraction and OCR processing
    - Database connectivity (SQL, MongoDB)
    - Markdown processing
    """
    
    SUPPORTED_EXTENSIONS = {
        '.pdf', '.docx', '.txt', '.md', '.ppt', '.pptx', 
        '.xlsx', '.xls', '.csv', '.db', '.sqlite'
    }
    
    def __init__(self, extract_images: bool = True, extract_tables: bool = True):
        self.extract_images = extract_images
        self.extract_tables = extract_tables
        
    def load_document(self, file_path: Union[str, Path]) -> List[Document]:
        """Load document based on file extension with enhanced processing."""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        try:
            if extension == '.pdf':
                return self._load_pdf_enhanced(file_path)
            elif extension in ['.docx']:
                return self._load_docx_enhanced(file_path)
            elif extension == '.txt':
                return self._load_text(file_path)
            elif extension == '.md':
                return self._load_markdown(file_path)
            elif extension in ['.ppt', '.pptx']:
                return self._load_powerpoint(file_path)
            elif extension in ['.xlsx', '.xls']:
                return self._load_excel(file_path)
            elif extension == '.csv':
                return self._load_csv(file_path)
            elif extension in ['.db', '.sqlite']:
                return self._load_sqlite(file_path)
            else:
                raise ValueError(f"Unsupported file extension: {extension}")
                
        except Exception as e:
            log.error(f"Error loading document {file_path}: {str(e)}")
            raise DocumentPortalException(f"Failed to load document: {file_path}", e)
    
    def _load_pdf_enhanced(self, file_path: Path) -> List[Document]:
        """Enhanced PDF loading with table and image extraction."""
        documents = []
        
        try:
            # Standard text extraction
            loader = PyPDFLoader(str(file_path))
            base_docs = loader.load()
            
            # Enhanced processing with PyMuPDF
            pdf_doc = fitz.open(file_path)
            
            for page_num in range(len(pdf_doc)):
                page = pdf_doc[page_num]
                content_parts = []
                
                # Extract text
                text = page.get_text()
                if text.strip():
                    content_parts.append(f"Text Content:\n{text}")
                
                # Extract tables if enabled
                if self.extract_tables:
                    tables = self._extract_tables_from_pdf_page(page)
                    if tables:
                        content_parts.append(f"Tables:\n{tables}")
                
                # Extract images if enabled
                if self.extract_images:
                    images_text = self._extract_images_from_pdf_page(page, page_num)
                    if images_text:
                        content_parts.append(f"Images OCR:\n{images_text}")
                
                if content_parts:
                    combined_content = "\n\n".join(content_parts)
                    doc = Document(
                        page_content=combined_content,
                        metadata={
                            "source": str(file_path),
                            "page": page_num + 1,
                            "file_type": "pdf",
                            "has_tables": self.extract_tables,
                            "has_images": self.extract_images
                        }
                    )
                    documents.append(doc)
            
            pdf_doc.close()
            log.info(f"Enhanced PDF loaded: {file_path}, pages: {len(documents)}")
            return documents
            
        except Exception as e:
            log.error(f"Error in enhanced PDF loading: {str(e)}")
            raise DocumentPortalException(f"Enhanced PDF loading failed: {file_path}", e)
    
    def _extract_tables_from_pdf_page(self, page) -> str:
        """Extract tables from PDF page using PyMuPDF."""
        try:
            tables = page.find_tables()
            table_texts = []
            
            for table in tables:
                df = table.to_pandas()
                table_text = df.to_string(index=False)
                table_texts.append(table_text)
            
            return "\n\n".join(table_texts) if table_texts else ""
        except Exception as e:
            log.warning(f"Table extraction failed: {str(e)}")
            return ""
    
    def _extract_images_from_pdf_page(self, page, page_num: int) -> str:
        """Extract and OCR images from PDF page."""
        try:
            image_list = page.get_images()
            ocr_texts = []
            
            for img_index, img in enumerate(image_list):
                xref = img[0]
                pix = fitz.Pixmap(page.parent, xref)
                
                if pix.n - pix.alpha < 4:  # GRAY or RGB
                    img_data = pix.tobytes("png")
                    img = Image.open(io.BytesIO(img_data))
                    
                    # Perform OCR
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text.strip():
                        ocr_texts.append(f"Image {img_index + 1}: {ocr_text.strip()}")
                
                pix = None
            
            return "\n\n".join(ocr_texts) if ocr_texts else ""
        except Exception as e:
            log.warning(f"Image extraction failed: {str(e)}")
            return ""
    
    def _load_docx_enhanced(self, file_path: Path) -> List[Document]:
        """Enhanced DOCX loading."""
        try:
            loader = Docx2txtLoader(str(file_path))
            docs = loader.load()
            
            # Enhance metadata
            for doc in docs:
                doc.metadata.update({
                    "file_type": "docx",
                    "source": str(file_path)
                })
            
            return docs
        except Exception as e:
            raise DocumentPortalException(f"DOCX loading failed: {file_path}", e)
    
    def _load_text(self, file_path: Path) -> List[Document]:
        """Load plain text files."""
        try:
            loader = TextLoader(str(file_path), encoding="utf-8")
            docs = loader.load()
            
            for doc in docs:
                doc.metadata.update({
                    "file_type": "txt",
                    "source": str(file_path)
                })
            
            return docs
        except Exception as e:
            raise DocumentPortalException(f"Text loading failed: {file_path}", e)
    
    def _load_markdown(self, file_path: Path) -> List[Document]:
        """Load and process Markdown files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # Convert markdown to HTML then extract text
            html = markdown.markdown(md_content)
            
            doc = Document(
                page_content=md_content,
                metadata={
                    "source": str(file_path),
                    "file_type": "markdown",
                    "html_version": html
                }
            )
            
            return [doc]
        except Exception as e:
            raise DocumentPortalException(f"Markdown loading failed: {file_path}", e)
    
    def _load_powerpoint(self, file_path: Path) -> List[Document]:
        """Load PowerPoint presentations with slide-by-slide processing."""
        try:
            prs = Presentation(file_path)
            documents = []
            
            for slide_num, slide in enumerate(prs.slides):
                content_parts = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text.strip())
                
                # Extract tables if present
                if self.extract_tables:
                    for shape in slide.shapes:
                        if shape.has_table:
                            table_text = self._extract_ppt_table(shape.table)
                            if table_text:
                                content_parts.append(f"Table:\n{table_text}")
                
                if content_parts:
                    slide_content = "\n\n".join(content_parts)
                    doc = Document(
                        page_content=slide_content,
                        metadata={
                            "source": str(file_path),
                            "slide": slide_num + 1,
                            "file_type": "powerpoint"
                        }
                    )
                    documents.append(doc)
            
            log.info(f"PowerPoint loaded: {file_path}, slides: {len(documents)}")
            return documents
            
        except Exception as e:
            raise DocumentPortalException(f"PowerPoint loading failed: {file_path}", e)
    
    def _extract_ppt_table(self, table) -> str:
        """Extract table data from PowerPoint table."""
        try:
            rows = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                rows.append(row_data)
            
            df = pd.DataFrame(rows[1:], columns=rows[0] if rows else [])
            return df.to_string(index=False)
        except Exception as e:
            log.warning(f"PPT table extraction failed: {str(e)}")
            return ""
    
    def _load_excel(self, file_path: Path) -> List[Document]:
        """Load Excel files with sheet-by-sheet processing."""
        try:
            workbook = load_workbook(file_path, data_only=True)
            documents = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Convert to pandas DataFrame
                data = []
                for row in sheet.iter_rows(values_only=True):
                    data.append(row)
                
                if data:
                    df = pd.DataFrame(data[1:], columns=data[0] if data else [])
                    sheet_content = df.to_string(index=False)
                    
                    doc = Document(
                        page_content=sheet_content,
                        metadata={
                            "source": str(file_path),
                            "sheet_name": sheet_name,
                            "file_type": "excel",
                            "rows": len(df),
                            "columns": len(df.columns)
                        }
                    )
                    documents.append(doc)
            
            log.info(f"Excel loaded: {file_path}, sheets: {len(documents)}")
            return documents
            
        except Exception as e:
            raise DocumentPortalException(f"Excel loading failed: {file_path}", e)
    
    def _load_csv(self, file_path: Path) -> List[Document]:
        """Load CSV files."""
        try:
            loader = CSVLoader(str(file_path))
            docs = loader.load()
            
            # Also create a summary document with statistics
            df = pd.read_csv(file_path)
            summary = f"""
CSV File Summary:
- Rows: {len(df)}
- Columns: {len(df.columns)}
- Column Names: {', '.join(df.columns.tolist())}
- Data Types: {df.dtypes.to_string()}

Sample Data:
{df.head().to_string()}
            """
            
            summary_doc = Document(
                page_content=summary,
                metadata={
                    "source": str(file_path),
                    "file_type": "csv_summary",
                    "rows": len(df),
                    "columns": len(df.columns)
                }
            )
            
            docs.append(summary_doc)
            return docs
            
        except Exception as e:
            raise DocumentPortalException(f"CSV loading failed: {file_path}", e)
    
    def _load_sqlite(self, file_path: Path) -> List[Document]:
        """Load SQLite database with table schemas and sample data."""
        try:
            conn = sqlite3.connect(file_path)
            documents = []
            
            # Get all table names
            cursor = conn.execute("SELECT name FROM sqlite_master WHERE type='table';")
            tables = cursor.fetchall()
            
            for table_name, in tables:
                # Get table schema
                schema_cursor = conn.execute(f"PRAGMA table_info({table_name});")
                schema = schema_cursor.fetchall()
                
                # Get sample data
                sample_cursor = conn.execute(f"SELECT * FROM {table_name} LIMIT 10;")
                sample_data = sample_cursor.fetchall()
                
                # Create document content
                schema_text = "Schema:\n" + "\n".join([f"- {col[1]} ({col[2]})" for col in schema])
                
                if sample_data:
                    df = pd.DataFrame(sample_data, columns=[col[1] for col in schema])
                    sample_text = f"\nSample Data:\n{df.to_string(index=False)}"
                else:
                    sample_text = "\nNo sample data available."
                
                content = f"Table: {table_name}\n{schema_text}{sample_text}"
                
                doc = Document(
                    page_content=content,
                    metadata={
                        "source": str(file_path),
                        "table_name": table_name,
                        "file_type": "sqlite",
                        "schema_columns": len(schema),
                        "sample_rows": len(sample_data)
                    }
                )
                documents.append(doc)
            
            conn.close()
            log.info(f"SQLite loaded: {file_path}, tables: {len(documents)}")
            return documents
            
        except Exception as e:
            raise DocumentPortalException(f"SQLite loading failed: {file_path}", e)
    
    def load_from_database_url(self, db_url: str, query: str) -> List[Document]:
        """Load data from database using connection URL and SQL query."""
        try:
            engine = create_engine(db_url)
            df = pd.read_sql(query, engine)
            
            content = f"Query: {query}\n\nResults:\n{df.to_string(index=False)}"
            
            doc = Document(
                page_content=content,
                metadata={
                    "source": db_url,
                    "query": query,
                    "file_type": "database_query",
                    "rows": len(df),
                    "columns": len(df.columns)
                }
            )
            
            return [doc]
            
        except Exception as e:
            raise DocumentPortalException(f"Database query failed: {db_url}", e)
    
    def load_from_mongodb(self, connection_string: str, database: str, collection: str, query: Dict = None) -> List[Document]:
        """Load data from MongoDB."""
        try:
            client = pymongo.MongoClient(connection_string)
            db = client[database]
            coll = db[collection]
            
            query = query or {}
            documents = []
            
            for doc in coll.find(query).limit(100):  # Limit for performance
                content = str(doc)
                
                lang_doc = Document(
                    page_content=content,
                    metadata={
                        "source": f"{connection_string}/{database}/{collection}",
                        "file_type": "mongodb",
                        "document_id": str(doc.get("_id", ""))
                    }
                )
                documents.append(lang_doc)
            
            client.close()
            return documents
            
        except Exception as e:
            raise DocumentPortalException(f"MongoDB loading failed: {connection_string}", e)


def load_documents_enhanced(paths: List[Path], **loader_kwargs) -> List[Document]:
    """Enhanced document loading function supporting all formats."""
    loader = EnhancedDocumentLoader(**loader_kwargs)
    all_documents = []
    
    for path in paths:
        try:
            docs = loader.load_document(path)
            all_documents.extend(docs)
        except Exception as e:
            log.error(f"Failed to load {path}: {str(e)}")
            continue
    
    log.info(f"Total documents loaded: {len(all_documents)}")
    return all_documents
